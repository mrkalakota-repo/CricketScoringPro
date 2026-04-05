"""
Build Inningsly — Executive Presentation
Run: python3 scripts/build_presentation.py
Output: docs/Inningsly_Executive.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Brand colours ──────────────────────────────────────────────────────────────
GREEN       = RGBColor(0x1B, 0x6B, 0x28)   # primary
GREEN_LIGHT = RGBColor(0xEA, 0xF7, 0xEB)   # light bg
GREEN_MID   = RGBColor(0x5E, 0xBD, 0x6A)   # dark-mode primary (accent)
ORANGE      = RGBColor(0xE6, 0x51, 0x00)   # secondary / cricket ball
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x1A, 0x1A, 0x1A)
GREY        = RGBColor(0x55, 0x55, 0x55)
LIGHT_GREY  = RGBColor(0xF4, 0xF4, 0xF4)
SLATE       = RGBColor(0x37, 0x47, 0x4F)
TEAL        = RGBColor(0x00, 0x69, 0x5C)

W = Inches(13.33)   # widescreen 16:9 width
H = Inches(7.5)     # widescreen 16:9 height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helper utilities ───────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def add_bullet_box(slide, bullets, x, y, w, h,
                   size=14, color=DARK, bullet_color=GREEN, header=None, header_size=15):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if header:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = header
        run.font.size = Pt(header_size)
        run.font.bold = True
        run.font.color.rgb = bullet_color
        first = False
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"  •  {b}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return txBox


def section_header(slide, label, color=GREEN):
    """Thin coloured bar under slide title."""
    add_rect(slide, Inches(0.5), Inches(1.55), Inches(12.33), Pt(3), fill=color)
    add_text(slide, label.upper(),
             Inches(0.5), Inches(0.28), Inches(12), Inches(0.8),
             size=9, bold=True, color=color, align=PP_ALIGN.LEFT)


def slide_title(slide, title, subtitle=None, title_color=DARK):
    add_text(slide, title,
             Inches(0.5), Inches(0.45), Inches(12.3), Inches(0.95),
             size=30, bold=True, color=title_color, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(1.22), Inches(12.3), Inches(0.45),
                 size=15, bold=False, color=GREY, align=PP_ALIGN.LEFT)


def card(slide, x, y, w, h, bg=WHITE, border=None):
    r = add_rect(slide, x, y, w, h, fill=bg,
                 line=border or RGBColor(0xE0, 0xE0, 0xE0),
                 line_width=Pt(0.75))
    return r


def icon_badge(slide, x, y, size_in, emoji_or_letter, bg=GREEN, fg=WHITE, font_size=20):
    sz = Inches(size_in)
    # circle approximation via rounded rect
    s = slide.shapes.add_shape(5, x, y, sz, sz)   # 5 = rounded rect
    s.adjustments[0] = 0.5
    s.fill.solid()
    s.fill.fore_color.rgb = bg
    s.line.fill.background()
    tf = s.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = emoji_or_letter
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = fg
    return s


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Cover
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)

# Full green left panel
add_rect(sl, 0, 0, Inches(5.4), H, fill=GREEN)
# Subtle diagonal accent strip
add_rect(sl, Inches(4.9), 0, Inches(0.6), H, fill=GREEN_MID)
# Right white area
add_rect(sl, Inches(5.5), 0, Inches(7.83), H, fill=WHITE)

# Cricket ball circle (decorative)
ball = sl.shapes.add_shape(9, Inches(0.4), Inches(4.8), Inches(1.6), Inches(1.6))  # 9=oval
ball.fill.solid(); ball.fill.fore_color.rgb = ORANGE
ball.line.fill.background()

ball2 = sl.shapes.add_shape(9, Inches(3.2), Inches(0.2), Inches(0.9), Inches(0.9))
ball2.fill.solid(); ball2.fill.fore_color.rgb = RGBColor(0xFF, 0x99, 0x44)
ball2.line.fill.background()

# App name (left panel)
add_text(sl, "🏏", Inches(0.55), Inches(1.0), Inches(4.5), Inches(0.9),
         size=48, color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "Inningsly",
         Inches(0.55), Inches(1.85), Inches(4.5), Inches(1.9),
         size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "Executive Overview",
         Inches(0.55), Inches(3.7), Inches(4.5), Inches(0.5),
         size=15, color=RGBColor(0xA8, 0xDB, 0xAB), align=PP_ALIGN.LEFT)
add_text(sl, "March 2026",
         Inches(0.55), Inches(4.15), Inches(4.5), Inches(0.4),
         size=12, color=RGBColor(0xA8, 0xDB, 0xAB), align=PP_ALIGN.LEFT)

# Tagline (right panel)
add_text(sl, "Score. Track.\nConnect.",
         Inches(5.9), Inches(1.5), Inches(6.8), Inches(2.0),
         size=42, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
add_text(sl,
         "A professional-grade cricket scoring platform\n"
         "for Android, iOS and Web — built for gully cricket,\n"
         "powered by real-time cloud technology.",
         Inches(5.9), Inches(3.5), Inches(6.8), Inches(1.5),
         size=14, color=GREY, align=PP_ALIGN.LEFT)

# Key stats row
stats = [("5", "Tabs"), ("4", "Formats"), ("50 mi", "Live Range"), ("∞", "Offline")]
for i, (val, lbl) in enumerate(stats):
    bx = Inches(5.9 + i * 1.75)
    add_rect(sl, bx, Inches(5.5), Inches(1.55), Inches(1.5),
             fill=GREEN_LIGHT, line=RGBColor(0xC0, 0xDC, 0xC2), line_width=Pt(0.5))
    add_text(sl, val, bx, Inches(5.55), Inches(1.55), Inches(0.7),
             size=24, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, bx, Inches(6.2), Inches(1.55), Inches(0.55),
             size=10, color=GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, Inches(1.65), fill=GREEN)
add_text(sl, "The Problem",
         Inches(0.6), Inches(0.35), Inches(11), Inches(1.0),
         size=30, bold=True, color=WHITE)
add_text(sl, "Why gully cricket needs a modern scoring app",
         Inches(0.6), Inches(1.0), Inches(11), Inches(0.5),
         size=13, color=RGBColor(0xA8, 0xDB, 0xAB))

problems = [
    ("📝", "Paper Scorecards", "Easy to lose, smudge, or miscalculate — no digital record survives the match."),
    ("📵", "No Offline Play", "Existing apps require constant internet. Most gully grounds have poor connectivity."),
    ("🔗", "Isolated Data", "Scorers can't share live scores with spectators watching from afar."),
    ("⚙️", "Oversimplified Tools", "Generic apps don't handle all formats, dismissal types, or extras correctly."),
    ("👥", "No Team Identity", "No persistent team/player records means starting from scratch every match."),
    ("🏆", "No Leagues", "Amateur teams have no lightweight way to run tournaments with standings."),
]

cols = 3
for i, (ico, title, desc) in enumerate(problems):
    col = i % cols
    row = i // cols
    cx = Inches(0.45 + col * 4.25)
    cy = Inches(1.85 + row * 2.5)
    card(sl, cx, cy, Inches(4.0), Inches(2.25), bg=WHITE)
    add_rect(sl, cx, cy, Inches(4.0), Inches(0.06), fill=ORANGE)
    add_text(sl, ico, cx + Inches(0.15), cy + Inches(0.12), Inches(0.6), Inches(0.55),
             size=22, align=PP_ALIGN.LEFT)
    add_text(sl, title, cx + Inches(0.75), cy + Inches(0.18), Inches(3.1), Inches(0.45),
             size=13, bold=True, color=DARK)
    add_text(sl, desc, cx + Inches(0.15), cy + Inches(0.65), Inches(3.7), Inches(1.4),
             size=11, color=GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Our Solution
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, Inches(1.65), fill=GREEN)
add_text(sl, "Our Solution",
         Inches(0.6), Inches(0.35), Inches(11), Inches(1.0),
         size=30, bold=True, color=WHITE)
add_text(sl, "Inningsly — one app, every format, everywhere",
         Inches(0.6), Inches(1.0), Inches(11), Inches(0.5),
         size=13, color=RGBColor(0xA8, 0xDB, 0xAB))

# Left column — description
add_text(sl,
         "Inningsly is a professional-grade, "
         "offline-first scoring platform built for street and community cricket. "
         "It handles every format from T20 to Test, stores complete match history, "
         "and broadcasts live scores to spectators within 50 miles — "
         "all without requiring a backend subscription or user accounts.",
         Inches(0.5), Inches(1.85), Inches(5.5), Inches(2.2),
         size=13, color=DARK)

pillars = [
    (GREEN,  "Offline First",     "Full feature parity with zero internet. SQLite on-device."),
    (ORANGE, "All Formats",       "T20, ODI, Test (follow-on), Custom — ICC-accurate rules."),
    (TEAL,   "Real-Time Cloud",   "Live scores, team chat, delegate access via Supabase."),
    (SLATE,  "Team Management",   "Persistent teams, rosters, leagues, career stats."),
]
for i, (clr, ttl, desc) in enumerate(pillars):
    cy = Inches(1.85 + i * 1.33)
    add_rect(sl, Inches(0.5), cy + Inches(0.15), Inches(0.22), Inches(0.9), fill=clr)
    add_text(sl, ttl, Inches(0.85), cy + Inches(0.15), Inches(5.2), Inches(0.4),
             size=13, bold=True, color=clr)
    add_text(sl, desc, Inches(0.85), cy + Inches(0.52), Inches(5.2), Inches(0.5),
             size=11, color=GREY)

# Right column — feature grid
features = [
    "Ball-by-ball scoring with undo",
    "All dismissal types + free hit",
    "Extras: wide, no ball, bye, leg bye",
    "Partnerships + fall of wickets",
    "Full scorecard (batting + bowling)",
    "Career stats (avg, SR, economy, best)",
    "Leagues with standings + fixtures",
    "Admin PIN (SHA-256) per team",
    "Proximity team discovery (50 mi)",
    "Dark mode · MD3 design system",
]
add_rect(sl, Inches(6.2), Inches(1.85), Inches(6.6), Inches(5.4), fill=GREEN_LIGHT,
         line=RGBColor(0xC0, 0xDC, 0xC2), line_width=Pt(0.5))
add_text(sl, "FEATURE CHECKLIST",
         Inches(6.45), Inches(2.0), Inches(6.0), Inches(0.4),
         size=9, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
for i, feat in enumerate(features):
    cy = Inches(2.4 + i * 0.46)
    add_rect(sl, Inches(6.45), cy + Inches(0.08), Inches(0.22), Inches(0.22),
             fill=GREEN)
    add_text(sl, "✓", Inches(6.44), cy + Inches(0.03), Inches(0.3), Inches(0.35),
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, feat, Inches(6.8), cy, Inches(5.7), Inches(0.4),
             size=11, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Technology Stack
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, Inches(1.65), fill=SLATE)
add_text(sl, "Technology Stack",
         Inches(0.6), Inches(0.35), Inches(11), Inches(1.0),
         size=30, bold=True, color=WHITE)
add_text(sl, "Modern, proven, open-source — zero licensing cost",
         Inches(0.6), Inches(1.0), Inches(11), Inches(0.5),
         size=13, color=RGBColor(0xB0, 0xBE, 0xC5))

tech = [
    # (category, name, why, color)
    ("Framework",     "React Native\n+ Expo SDK 54",
     "Single codebase → Android, iOS, Web.\nExpo Go for instant device testing;\nEAS Build for app store deployment.",
     RGBColor(0x61, 0xDB, 0xFB)),
    ("Language",      "TypeScript\n(strict mode)",
     "Type-safe data models enforced at\ncompile time — catches cricket rule\nbugs before they ship.",
     RGBColor(0x3B, 0x82, 0xC4)),
    ("Navigation",    "Expo Router v6",
     "File-based routing with deep links.\nTab navigator + modal stack.\nURL-addressable screens.",
     RGBColor(0xFF, 0xFF, 0xFF)),
    ("UI Library",    "React Native\nPaper (MD3)",
     "Material Design 3 components.\nLight + dark theme tokens.\nWeb-compatible dialogs.",
     RGBColor(0xFF, 0x72, 0x43)),
    ("State",         "Zustand",
     "Lightweight global state with\nno boilerplate. Stores map 1:1\nto domain areas.",
     RGBColor(0xFF, 0xA0, 0x00)),
    ("Persistence",   "expo-sqlite\n/ localStorage",
     "SQLite on native; localStorage on\nweb via Metro .web.ts resolution.\nFully offline-first.",
     RGBColor(0x66, 0xBB, 0x6A)),
    ("Cloud",         "Supabase\n(PostgreSQL)",
     "Real-time live scores, team chat,\ndelegate codes. Row-level security.\nOptional — app works without it.",
     RGBColor(0x3E, 0xCF, 0x8E)),
    ("Testing",       "Jest (150+ tests)",
     "Pure engine unit tests + 99\nfunctional end-to-end scenarios.\nAll cricket rules validated.",
     RGBColor(0xC6, 0x78, 0xDD)),
]

cols = 4
for i, (cat, name, why, clr) in enumerate(tech):
    col = i % cols
    row = i // cols
    cx = Inches(0.35 + col * 3.2)
    cy = Inches(1.85 + row * 2.65)
    # Card
    add_rect(sl, cx, cy, Inches(3.0), Inches(2.45),
             fill=RGBColor(0x26, 0x32, 0x38),
             line=clr, line_width=Pt(1.5))
    # Top colour stripe
    add_rect(sl, cx, cy, Inches(3.0), Inches(0.08), fill=clr)
    # Category label
    add_text(sl, cat.upper(), cx + Inches(0.15), cy + Inches(0.14),
             Inches(2.7), Inches(0.3),
             size=8, bold=True, color=clr)
    # Tech name
    add_text(sl, name, cx + Inches(0.15), cy + Inches(0.42),
             Inches(2.7), Inches(0.75),
             size=13, bold=True, color=WHITE)
    # Why
    add_text(sl, why, cx + Inches(0.15), cy + Inches(1.18),
             Inches(2.7), Inches(1.15),
             size=9.5, color=RGBColor(0xB0, 0xBE, 0xC5))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Architecture
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, Inches(1.65), fill=GREEN)
add_text(sl, "Architecture Overview",
         Inches(0.6), Inches(0.35), Inches(11), Inches(1.0),
         size=30, bold=True, color=WHITE)
add_text(sl, "Layered, offline-first design with optional cloud augmentation",
         Inches(0.6), Inches(1.0), Inches(11), Inches(0.5),
         size=13, color=RGBColor(0xA8, 0xDB, 0xAB))

# Layer boxes (left-to-right stack)
layers = [
    (GREEN_LIGHT,                   "PRESENTATION",  "Expo Router screens\nReact Native Paper\nDark / Light MD3 theme"),
    (RGBColor(0xD1, 0xF0, 0xD3),    "STATE",         "Zustand stores\n(match, team, league,\nchat, live scores)"),
    (RGBColor(0xA8, 0xDB, 0xAB),    "ENGINE",        "Pure TypeScript\nImmutable MatchEngine\n150+ unit tests"),
    (RGBColor(0x7C, 0xBF, 0x81),    "PERSISTENCE",   "expo-sqlite (native)\nlocalStorage (web)\nSame repo interface"),
    (GREEN,                          "CLOUD",         "Supabase (optional)\nLive scores · Chat\nDelegate codes"),
]

arrow_x = Inches(0.45)
for i, (bg, lbl, desc) in enumerate(layers):
    bx = Inches(0.35 + i * 2.55)
    add_rect(sl, bx, Inches(1.85), Inches(2.35), Inches(3.8), fill=bg,
             line=GREEN, line_width=Pt(0.5))
    add_text(sl, lbl, bx + Inches(0.1), Inches(1.9), Inches(2.15), Inches(0.4),
             size=9, bold=True, color=GREEN if i < 4 else WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, desc, bx + Inches(0.1), Inches(2.35), Inches(2.15), Inches(1.8),
             size=11, color=DARK if i < 4 else WHITE, align=PP_ALIGN.CENTER)
    # Arrow (except after last)
    if i < len(layers) - 1:
        add_text(sl, "→", Inches(0.35 + i * 2.55 + 2.35), Inches(2.5),
                 Inches(0.25), Inches(0.6), size=18, bold=True, color=GREEN,
                 align=PP_ALIGN.CENTER)

# Key design principles (below layers)
add_rect(sl, Inches(0.35), Inches(6.0), Inches(12.6), Inches(1.3),
         fill=RGBColor(0xF0, 0xF8, 0xF0), line=RGBColor(0xC0, 0xDC, 0xC2), line_width=Pt(0.5))

principles = [
    ("Immutable Engine", "MatchEngine returns a new instance on every ball — undo is O(1)"),
    ("Platform Repos",   "Metro resolves .web.ts — SQLite never bundled for browser"),
    ("Cloud Optional",   "isCloudEnabled guards every Supabase call — full offline parity"),
    ("Store as Gateway", "UI never calls repos directly — Zustand stores are the only entry point"),
]
for i, (ttl, desc) in enumerate(principles):
    bx = Inches(0.55 + i * 3.1)
    add_text(sl, ttl, bx, Inches(6.08), Inches(2.9), Inches(0.38),
             size=10, bold=True, color=GREEN)
    add_text(sl, desc, bx, Inches(6.44), Inches(2.9), Inches(0.7),
             size=9, color=GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Key Features: Scoring Engine
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, Inches(1.65), fill=ORANGE)
add_text(sl, "Ball-by-Ball Scoring Engine",
         Inches(0.6), Inches(0.35), Inches(11), Inches(1.0),
         size=30, bold=True, color=WHITE)
add_text(sl, "ICC-accurate rules, tested to 150+ scenarios",
         Inches(0.6), Inches(1.0), Inches(11), Inches(0.5),
         size=13, color=RGBColor(0xFF, 0xCC, 0xAA))

# Scoring panel mockup (left)
add_rect(sl, Inches(0.4), Inches(1.85), Inches(5.5), Inches(5.4),
         fill=RGBColor(0x1A, 0x1A, 0x2E), line=RGBColor(0x44, 0x44, 0x66), line_width=Pt(1))
# Score line
add_text(sl, "THB  142 / 6  (17.3 ov)", Inches(0.55), Inches(2.0),
         Inches(5.2), Inches(0.5), size=13, bold=True, color=WHITE)
add_text(sl, "CRR: 8.11   ▪   RRR: 12.40   ▪   TGT: 186",
         Inches(0.55), Inches(2.45), Inches(5.2), Inches(0.35),
         size=9, color=RGBColor(0xB0, 0xB0, 0xC8))
# Batters
add_rect(sl, Inches(0.55), Inches(2.85), Inches(5.2), Inches(0.04),
         fill=RGBColor(0x44, 0x44, 0x66))
add_text(sl, "🏏  Ravi Kumar*      34(28)   4s:2  6s:1",
         Inches(0.55), Inches(2.92), Inches(5.2), Inches(0.38),
         size=10, color=RGBColor(0xFF, 0xD7, 0x00))
add_text(sl, "     Arjun Singh      18(15)",
         Inches(0.55), Inches(3.28), Inches(5.2), Inches(0.35),
         size=10, color=WHITE)
add_text(sl, "🎯  Deepak Rao   2-0-18-1  (4.3 ov)",
         Inches(0.55), Inches(3.65), Inches(5.2), Inches(0.38),
         size=10, color=RGBColor(0x80, 0xFF, 0xB4))
add_text(sl, "This over:   ·   1   4   W   ·   2",
         Inches(0.55), Inches(4.0), Inches(5.2), Inches(0.38),
         size=10, color=RGBColor(0xB0, 0xB0, 0xC8), italic=True)
# Button grid
btn_labels = ["0","1","2","3","4","6","W","↩"]
btn_colors_bg = [
    RGBColor(0x2A,0x2A,0x3E), RGBColor(0x2A,0x2A,0x3E), RGBColor(0x2A,0x2A,0x3E), RGBColor(0x2A,0x2A,0x3E),
    RGBColor(0x1B,0x3A,0x1B), RGBColor(0x1B,0x3A,0x1B), RGBColor(0x5C,0x10,0x10), RGBColor(0x33,0x27,0x00),
]
btn_colors_fg = [WHITE, WHITE, WHITE, WHITE,
                 RGBColor(0x80,0xFF,0x80), RGBColor(0x80,0xFF,0x80),
                 RGBColor(0xFF,0x70,0x70), RGBColor(0xFF,0xCC,0x00)]
for i, (lbl, bg, fg) in enumerate(zip(btn_labels, btn_colors_bg, btn_colors_fg)):
    bx = Inches(0.55 + (i % 4) * 1.3)
    by = Inches(4.5 + (i // 4) * 0.7)
    add_rect(sl, bx, by, Inches(1.18), Inches(0.58), fill=bg,
             line=RGBColor(0x55,0x55,0x77), line_width=Pt(0.75))
    add_text(sl, lbl, bx, by + Inches(0.08), Inches(1.18), Inches(0.45),
             size=14, bold=True, color=fg, align=PP_ALIGN.CENTER)
# Extras row
for j, ext in enumerate(["Wide","No Ball","Bye","Leg Bye"]):
    bx = Inches(0.55 + j * 1.3)
    add_rect(sl, bx, Inches(5.85), Inches(1.18), Inches(0.42),
             fill=RGBColor(0x22,0x22,0x44),
             line=RGBColor(0x55,0x55,0x77), line_width=Pt(0.5))
    add_text(sl, ext, bx, Inches(5.87), Inches(1.18), Inches(0.38),
             size=8, color=RGBColor(0xCC,0xCC,0xFF), align=PP_ALIGN.CENTER)

# Right: feature bullets
rules = [
    ("All 12 Dismissal Types",  "Bowled, caught, LBW, run out, stumped, hit wicket,\nretired hurt/out, obstructing field, timed out, hit twice"),
    ("Accurate Strike Rotation", "Odd-run rotation applied before end-of-over swap;\nno-ball uses runs off bat only (not penalty run)"),
    ("Free Hit Tracking",        "Wide/NB triggers free hit on next delivery;\nonly run-out + retired-out valid — all others blocked"),
    ("Undo Stack",               "Every ball snapshots innings state; unlimited undos;\npersisted in match JSON (survives app restart)"),
    ("Extras Breakdown",         "Wides, no balls, byes, leg byes tracked separately;\ncontribute correctly to bowler and innings totals"),
    ("Powerplay Zones",          "T20: overs 1–6 mandatory;\nODI: 3 powerplay blocks; Test: none"),
    ("Maiden Detection",         "Over marked maiden if 0 runs in 6 legal deliveries\n(regardless of wickets in the over)"),
]
for i, (ttl, desc) in enumerate(rules):
    cy = Inches(1.85 + i * 0.75)
    add_rect(sl, Inches(6.2), cy + Inches(0.1), Inches(0.12), Inches(0.45), fill=ORANGE)
    add_text(sl, ttl, Inches(6.45), cy + Inches(0.08), Inches(6.4), Inches(0.35),
             size=11, bold=True, color=DARK)
    add_text(sl, desc, Inches(6.45), cy + Inches(0.4), Inches(6.4), Inches(0.38),
             size=9, color=GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Team & Player Management
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, Inches(1.65), fill=GREEN)
add_text(sl, "Team & Player Management",
         Inches(0.6), Inches(0.35), Inches(11), Inches(1.0),
         size=30, bold=True, color=WHITE)
add_text(sl, "Persistent rosters, roles, proximity discovery, and delegate access",
         Inches(0.6), Inches(1.0), Inches(11), Inches(0.5),
         size=13, color=RGBColor(0xA8, 0xDB, 0xAB))

# Three feature areas
areas = [
    {
        "title": "Team Creation & Discovery",
        "color": GREEN,
        "points": [
            "Create teams with name, short code (≤5 chars), optional admin PIN",
            "GPS coordinates captured at creation for proximity sorting",
            "My Teams section always visible first",
            "Nearby Teams: up to 10 within 50 miles (Haversine formula)",
            "Search across all teams by name",
            "Admin PIN hashed with SHA-256 — plaintext never stored",
        ],
    },
    {
        "title": "Player Profiles & Roles",
        "color": ORANGE,
        "points": [
            "Batting style: Right-hand / Left-hand",
            "Bowling style: 9 options (fast, medium, spin variants)",
            "Roles: Captain (C), Vice-Captain (VC), Wicket Keeper (WK), All-Rounder (AR)",
            "C and VC are mutually exclusive — enforced at entry",
            "Optional phone number as cross-team identity key",
            "6-character shareable player code for profile lookup",
        ],
    },
    {
        "title": "Delegate Access",
        "color": TEAL,
        "points": [
            "Owner generates a 6-char one-time code (10-min TTL)",
            "Code stored in Supabase — single-use, deleted on redeem",
            "Another device enters code → granted editor access",
            "Delegate access stored locally in user_prefs",
            "No need to share admin PIN",
            "Requires cloud (Supabase) to be configured",
        ],
    },
]

for i, area in enumerate(areas):
    bx = Inches(0.35 + i * 4.3)
    clr = area["color"]
    add_rect(sl, bx, Inches(1.85), Inches(4.1), Inches(5.45),
             fill=WHITE, line=clr, line_width=Pt(1.5))
    add_rect(sl, bx, Inches(1.85), Inches(4.1), Inches(0.55), fill=clr)
    add_text(sl, area["title"],
             bx + Inches(0.15), Inches(1.93), Inches(3.8), Inches(0.42),
             size=11, bold=True, color=WHITE)
    for j, pt in enumerate(area["points"]):
        cy = Inches(2.55 + j * 0.75)
        add_rect(sl, bx + Inches(0.2), cy + Inches(0.12),
                 Inches(0.18), Inches(0.18), fill=clr)
        add_text(sl, pt, bx + Inches(0.5), cy, Inches(3.45), Inches(0.65),
                 size=10, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Cloud & Real-Time Features
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, Inches(1.65), fill=TEAL)
add_text(sl, "Cloud & Real-Time Features",
         Inches(0.6), Inches(0.35), Inches(11), Inches(1.0),
         size=30, bold=True, color=WHITE)
add_text(sl, "Optional Supabase integration — app works fully offline without it",
         Inches(0.6), Inches(1.0), Inches(11), Inches(0.5),
         size=13, color=RGBColor(0xA8, 0xDB, 0xAB))

# Live Scores flow
add_rect(sl, Inches(0.35), Inches(1.85), Inches(12.6), Inches(2.1),
         fill=RGBColor(0xF0, 0xFD, 0xFA), line=TEAL, line_width=Pt(1))
add_text(sl, "📡  LIVE MATCH SCORES — 50-MILE BROADCAST",
         Inches(0.55), Inches(1.95), Inches(12.0), Inches(0.4),
         size=11, bold=True, color=TEAL)

flow_steps = [
    ("1", "Scorer records\na ball"),
    ("→", ""),
    ("2", "Match upserted\nto Supabase\nlive_matches"),
    ("→", ""),
    ("3", "Postgres Realtime\ntriggers channel\nevent"),
    ("→", ""),
    ("4", "Viewer device\nfetches nearby\nmatches (bbox)"),
    ("→", ""),
    ("5", "Home screen\nupdates in\nreal time"),
]
x0 = Inches(0.6)
for i, (num, label) in enumerate(flow_steps):
    bx = x0 + Inches(i * 1.4)
    if num == "→":
        add_text(sl, "→", bx, Inches(2.5), Inches(0.5), Inches(0.5),
                 size=16, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    else:
        add_rect(sl, bx, Inches(2.38), Inches(1.2), Inches(1.18),
                 fill=TEAL, line=TEAL, line_width=Pt(0))
        add_text(sl, num, bx, Inches(2.38), Inches(1.2), Inches(0.5),
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, label, bx + Inches(0.05), Inches(2.75), Inches(1.1), Inches(0.65),
                 size=8.5, color=DARK, align=PP_ALIGN.CENTER)

# Bottom cards
cloud_feats = [
    (TEAL,  "🗺️  Proximity Query",
     "Bounding-box query on lat/lon. Returns up to 20 matches\nupdated in the last 24 hours within 50-mile radius.\nNo PostGIS required — fast index-only scan."),
    (GREEN, "💬  Team Chat",
     "Per-team Supabase real-time channel. Last 100 messages\nloaded on open; new messages arrive via postgres_changes.\nNo user accounts — identity from player profile."),
    (ORANGE,"🔐  Security Model",
     "Anon key (public read/write) + Row Level Security.\nAdmin PIN is SHA-256 client-side only — never sent to cloud.\nDelegate codes single-use + time-limited (10 min TTL)."),
    (SLATE, "🌐  Graceful Degradation",
     "PGRST205 (table not found) silently ignored — cloud\nfeatures simply don't show until SQL is run in Supabase.\nAll four cloud tables optional independently."),
]
for i, (clr, ttl, desc) in enumerate(cloud_feats):
    bx = Inches(0.35 + i * 3.25)
    add_rect(sl, bx, Inches(4.1), Inches(3.1), Inches(3.2),
             fill=WHITE, line=clr, line_width=Pt(1.5))
    add_rect(sl, bx, Inches(4.1), Inches(3.1), Inches(0.5), fill=clr)
    add_text(sl, ttl, bx + Inches(0.12), Inches(4.17), Inches(2.86), Inches(0.38),
             size=10, bold=True, color=WHITE)
    add_text(sl, desc, bx + Inches(0.12), Inches(4.68), Inches(2.86), Inches(2.4),
             size=9.5, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Leagues & Stats
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, Inches(1.65), fill=RGBColor(0x4A, 0x14, 0x8C))  # deep purple
add_text(sl, "Leagues & Career Statistics",
         Inches(0.6), Inches(0.35), Inches(11), Inches(1.0),
         size=30, bold=True, color=WHITE)
add_text(sl, "Run tournaments, track career milestones, celebrate achievements",
         Inches(0.6), Inches(1.0), Inches(11), Inches(0.5),
         size=13, color=RGBColor(0xCE, 0x93, 0xD8))

# Leagues section
add_rect(sl, Inches(0.35), Inches(1.85), Inches(6.0), Inches(5.45),
         fill=WHITE, line=RGBColor(0x4A, 0x14, 0x8C), line_width=Pt(1.5))
add_rect(sl, Inches(0.35), Inches(1.85), Inches(6.0), Inches(0.55),
         fill=RGBColor(0x4A, 0x14, 0x8C))
add_text(sl, "🏅  LEAGUES",
         Inches(0.5), Inches(1.93), Inches(5.6), Inches(0.4),
         size=12, bold=True, color=WHITE)

# Standings mockup
add_text(sl, "#   Team            P   W   L   Pts",
         Inches(0.5), Inches(2.52), Inches(5.7), Inches(0.38),
         size=10, bold=True, color=RGBColor(0x4A, 0x14, 0x8C))
rows = [
    ("1", "Thunderbolts",  "4", "4", "0", "8"),
    ("2", "Royal Strikers","4", "3", "1", "6"),
    ("3", "Blaze XI",      "4", "1", "3", "2"),
    ("4", "Sky Hawks",     "4", "0", "4", "0"),
]
for i, (pos, team, p, w, l, pts) in enumerate(rows):
    cy = Inches(2.9 + i * 0.55)
    bg = RGBColor(0xF3, 0xE5, 0xF5) if i == 0 else WHITE
    add_rect(sl, Inches(0.45), cy, Inches(5.7), Inches(0.48), fill=bg)
    row_text = f"  {pos}   {team:<18}{p}   {w}   {l}   {pts}"
    add_text(sl, row_text, Inches(0.5), cy + Inches(0.07), Inches(5.6), Inches(0.38),
             size=10, color=DARK if i > 0 else RGBColor(0x4A, 0x14, 0x8C),
             bold=(i == 0))

league_pts = [
    "Create leagues with any number of teams",
    "Round-robin fixture generation",
    "Link real match results to fixtures",
    "Live standings update automatically",
    "Scheduled, completed + abandoned states",
]
for i, pt in enumerate(league_pts):
    add_rect(sl, Inches(0.55), Inches(5.05 + i * 0.38), Inches(0.14), Inches(0.14),
             fill=RGBColor(0x4A, 0x14, 0x8C))
    add_text(sl, pt, Inches(0.8), Inches(5.02 + i * 0.38),
             Inches(5.4), Inches(0.38), size=10, color=DARK)

# Stats section
add_rect(sl, Inches(6.7), Inches(1.85), Inches(6.25), Inches(5.45),
         fill=WHITE, line=GREEN, line_width=Pt(1.5))
add_rect(sl, Inches(6.7), Inches(1.85), Inches(6.25), Inches(0.55), fill=GREEN)
add_text(sl, "📊  CAREER STATISTICS",
         Inches(6.85), Inches(1.93), Inches(5.9), Inches(0.4),
         size=12, bold=True, color=WHITE)

bat_stats = [("Matches", "8"), ("Innings", "7"), ("Runs", "312"),
             ("Highest", "87"), ("Average", "52.0"), ("Strike Rate", "145.8"),
             ("50s / 100s", "2 / 0"), ("Fours / Sixes", "18 / 9")]
bowl_stats = [("Overs", "12"), ("Wickets", "8"), ("Runs", "96"),
              ("Economy", "8.0"), ("Average", "12.0"), ("Best Figures", "3/18")]

add_text(sl, "BATTING — Ravi Kumar", Inches(6.85), Inches(2.5),
         Inches(5.9), Inches(0.38), size=10, bold=True, color=GREEN)
for i, (lbl, val) in enumerate(bat_stats):
    col = i % 2; row2 = i // 2
    bx = Inches(6.85 + col * 3.0)
    cy = Inches(2.88 + row2 * 0.48)
    add_text(sl, f"{lbl}:", bx, cy, Inches(1.7), Inches(0.4), size=9.5, color=GREY)
    add_text(sl, val, bx + Inches(1.6), cy, Inches(1.35), Inches(0.4),
             size=9.5, bold=True, color=DARK)

add_text(sl, "BOWLING — Ravi Kumar", Inches(6.85), Inches(4.8),
         Inches(5.9), Inches(0.38), size=10, bold=True, color=GREEN)
for i, (lbl, val) in enumerate(bowl_stats):
    col = i % 2; row2 = i // 2
    bx = Inches(6.85 + col * 3.0)
    cy = Inches(5.18 + row2 * 0.48)
    add_text(sl, f"{lbl}:", bx, cy, Inches(1.7), Inches(0.4), size=9.5, color=GREY)
    add_text(sl, val, bx + Inches(1.6), cy, Inches(1.35), Inches(0.4),
             size=9.5, bold=True, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Platform & Deployment
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, Inches(1.65), fill=SLATE)
add_text(sl, "Platform & Deployment",
         Inches(0.6), Inches(0.35), Inches(11), Inches(1.0),
         size=30, bold=True, color=WHITE)
add_text(sl, "One codebase — Android, iOS, and mobile web",
         Inches(0.6), Inches(1.0), Inches(11), Inches(0.5),
         size=13, color=RGBColor(0xB0, 0xBE, 0xC5))

platforms = [
    ("🤖", "Android",
     "Native APK via EAS Build\nExpo Go for development\nSQLite persistence\nTarget: Android 8+",
     RGBColor(0x3D, 0xDC, 0x84)),
    ("🍎", "iOS",
     "Native IPA via EAS Build\nExpo Go for development\nSQLite persistence\nTarget: iOS 15+",
     RGBColor(0xA2, 0xAA, 0xAD)),
    ("🌐", "Mobile Web",
     "Progressive Web App\nRun: expo start --web\nlocalStorage persistence\nMD3 responsive layout",
     RGBColor(0x61, 0xDB, 0xFB)),
]
for i, (ico, name, desc, clr) in enumerate(platforms):
    bx = Inches(0.35 + i * 4.3)
    add_rect(sl, bx, Inches(1.85), Inches(4.1), Inches(3.0),
             fill=WHITE, line=clr, line_width=Pt(2))
    add_text(sl, ico, bx + Inches(0.2), Inches(2.0), Inches(1.0), Inches(0.7),
             size=32, align=PP_ALIGN.LEFT)
    add_text(sl, name, bx + Inches(1.1), Inches(2.08), Inches(2.8), Inches(0.55),
             size=18, bold=True, color=clr)
    add_text(sl, desc, bx + Inches(0.2), Inches(2.75), Inches(3.7), Inches(1.8),
             size=11, color=DARK)

# Deployment section
add_rect(sl, Inches(0.35), Inches(5.05), Inches(12.6), Inches(2.2),
         fill=RGBColor(0x27, 0x32, 0x38), line=RGBColor(0x44, 0x55, 0x60), line_width=Pt(0.5))
add_text(sl, "DEPLOYMENT PIPELINE",
         Inches(0.55), Inches(5.15), Inches(12.0), Inches(0.38),
         size=10, bold=True, color=RGBColor(0xB0, 0xBE, 0xC5))

deploy_steps = [
    ("Dev", "Expo Go on device\n(LAN, instant reload)"),
    ("Test", "Jest 150+ tests\n(engine rules)"),
    ("Build", "EAS Build\n(cloud compile)"),
    ("Distribute", "Internal / TestFlight\n/ Play Store"),
    ("Cloud", "Run supabase-setup.sql\nAdd .env creds"),
]
for i, (stage, desc) in enumerate(deploy_steps):
    bx = Inches(0.55 + i * 2.5)
    add_rect(sl, bx, Inches(5.6), Inches(2.25), Inches(1.45),
             fill=RGBColor(0x37, 0x47, 0x4F), line=RGBColor(0x55, 0x66, 0x70), line_width=Pt(0.75))
    add_text(sl, stage.upper(), bx + Inches(0.1), Inches(5.65),
             Inches(2.05), Inches(0.38), size=9, bold=True,
             color=RGBColor(0x80, 0xCB, 0xC4))
    add_text(sl, desc, bx + Inches(0.1), Inches(6.0),
             Inches(2.05), Inches(0.9), size=9.5, color=WHITE)
    if i < 4:
        add_text(sl, "→", Inches(0.55 + i * 2.5 + 2.25), Inches(6.0),
                 Inches(0.28), Inches(0.38), size=14, bold=True,
                 color=RGBColor(0x80, 0xCB, 0xC4), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Security & Privacy
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, Inches(1.65), fill=RGBColor(0x18, 0x3A, 0x28))
add_text(sl, "Security & Privacy",
         Inches(0.6), Inches(0.35), Inches(11), Inches(1.0),
         size=30, bold=True, color=WHITE)
add_text(sl, "Built with data minimisation and user trust at the core",
         Inches(0.6), Inches(1.0), Inches(11), Inches(0.5),
         size=13, color=RGBColor(0xA8, 0xDB, 0xAB))

security = [
    ("🔒", "Admin PIN (SHA-256)",
     "Team admin PINs hashed with SHA-256 via expo-crypto before storage.\n"
     "Plaintext never stored in SQLite or transmitted to cloud.\n"
     "Auth state held in-memory only — clears on app restart by design."),
    ("💾", "Local-First Data",
     "All match data, player records, and team info stored exclusively on-device\n"
     "(SQLite native / localStorage web). No automatic cloud upload of match data.\n"
     "Users own their data entirely."),
    ("🛡️", "Parameterised Queries",
     "All SQLite operations use parameterised statements — SQL injection is\n"
     "architecturally impossible. Input validation at store boundaries."),
    ("☁️", "Cloud Opt-In Only",
     "Supabase integration requires explicit credential setup in .env.\n"
     "isCloudEnabled guards every remote call — cloud is never silently enabled.\n"
     "PGRST205 errors (missing table) silently suppressed — graceful degradation."),
    ("🔑", "Delegate Codes",
     "Single-use, time-limited (10 min) codes for sharing team access.\n"
     "Stored in Supabase and deleted immediately on first successful claim.\n"
     "No persistent credentials shared — codes cannot be replayed."),
    ("📵", "No PII Transmission",
     "Player names and match data are never sent to Anthropic or any third party.\n"
     "Phone numbers are optional and stored locally only.\n"
     "Cloud tables contain only team metadata and live match scores."),
]
cols = 3
for i, (ico, ttl, desc) in enumerate(security):
    col = i % cols; row = i // cols
    bx = Inches(0.35 + col * 4.3)
    cy = Inches(1.85 + row * 2.7)
    add_rect(sl, bx, cy, Inches(4.1), Inches(2.5),
             fill=WHITE, line=GREEN, line_width=Pt(1))
    add_rect(sl, bx, cy, Inches(4.1), Inches(0.06), fill=GREEN)
    add_text(sl, ico, bx + Inches(0.15), cy + Inches(0.12), Inches(0.6), Inches(0.55),
             size=22)
    add_text(sl, ttl, bx + Inches(0.8), cy + Inches(0.18), Inches(3.2), Inches(0.42),
             size=12, bold=True, color=GREEN)
    add_text(sl, desc, bx + Inches(0.15), cy + Inches(0.7), Inches(3.8), Inches(1.65),
             size=9.5, color=GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Roadmap
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, Inches(1.65), fill=ORANGE)
add_text(sl, "Product Roadmap",
         Inches(0.6), Inches(0.35), Inches(11), Inches(1.0),
         size=30, bold=True, color=WHITE)
add_text(sl, "From foundation to community platform",
         Inches(0.6), Inches(1.0), Inches(11), Inches(0.5),
         size=13, color=RGBColor(0xFF, 0xCC, 0xAA))

phases = [
    {
        "phase": "Phase 1",
        "label": "COMPLETE ✓",
        "label_color": GREEN,
        "title": "Foundation",
        "color": GREEN,
        "items": [
            "T20, ODI, Test, Custom formats",
            "Ball-by-ball scoring + undo",
            "All dismissal types + free hit",
            "Team + player management",
            "Full scorecard + partnerships",
            "Career stats (batting + bowling)",
            "SQLite + localStorage persistence",
            "150+ unit + functional tests",
        ],
    },
    {
        "phase": "Phase 2",
        "label": "COMPLETE ✓",
        "label_color": GREEN,
        "title": "Cloud & Community",
        "color": TEAL,
        "items": [
            "Supabase cloud sync",
            "Live match scores (50-mile broadcast)",
            "Real-time team chat",
            "Delegate team access codes",
            "Leagues + standings",
            "Admin PIN (SHA-256)",
            "Dark mode + MD3 theme",
            "Expo Router navigation",
        ],
    },
    {
        "phase": "Phase 3",
        "label": "UPCOMING",
        "label_color": ORANGE,
        "title": "Growth",
        "color": ORANGE,
        "items": [
            "Push notifications (ball alerts)",
            "Match commentary sharing",
            "Photo / media attachments in chat",
            "Tournament bracket visualisation",
            "Landscape scorecard view",
            "Offline map with nearby matches",
            "Player performance trends",
            "EAS app store submission",
        ],
    },
]

for i, ph in enumerate(phases):
    bx = Inches(0.35 + i * 4.3)
    clr = ph["color"]
    lbl_clr = ph["label_color"]
    add_rect(sl, bx, Inches(1.85), Inches(4.1), Inches(5.45),
             fill=WHITE, line=clr, line_width=Pt(2))
    add_rect(sl, bx, Inches(1.85), Inches(4.1), Inches(0.9), fill=clr)
    add_text(sl, ph["phase"], bx + Inches(0.15), Inches(1.9),
             Inches(3.8), Inches(0.35), size=9, bold=True, color=WHITE)
    add_text(sl, ph["title"], bx + Inches(0.15), Inches(2.2),
             Inches(3.8), Inches(0.42), size=14, bold=True, color=WHITE)
    # Status badge
    add_rect(sl, bx + Inches(2.0), Inches(2.65), Inches(1.95), Inches(0.38),
             fill=WHITE if ph["label"] == "UPCOMING" else GREEN_LIGHT,
             line=lbl_clr, line_width=Pt(1))
    add_text(sl, ph["label"], bx + Inches(2.0), Inches(2.67),
             Inches(1.95), Inches(0.32), size=8, bold=True,
             color=lbl_clr, align=PP_ALIGN.CENTER)

    for j, item in enumerate(ph["items"]):
        cy = Inches(2.85 + j * 0.6)
        add_rect(sl, bx + Inches(0.2), cy + Inches(0.14),
                 Inches(0.14), Inches(0.14), fill=clr)
        add_text(sl, item, bx + Inches(0.45), cy,
                 Inches(3.5), Inches(0.55), size=10, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Summary / Thank You
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)

# Split layout
add_rect(sl, 0, 0, W, H, fill=GREEN)
add_rect(sl, Inches(5.5), 0, Inches(7.83), H, fill=WHITE)
add_rect(sl, Inches(5.3), 0, Inches(0.4), H, fill=GREEN_MID)

# Left — key messages
add_text(sl, "🏏", Inches(0.5), Inches(0.8), Inches(4.5), Inches(0.9),
         size=42, color=WHITE)
add_text(sl, "Inningsly",
         Inches(0.5), Inches(1.6), Inches(4.5), Inches(1.6),
         size=30, bold=True, color=WHITE)

msgs = [
    "✅  ICC-accurate scoring for every format",
    "✅  Works completely offline — no internet needed",
    "✅  Live scores broadcast to 50-mile radius",
    "✅  Teams, leagues, career stats, and chat",
    "✅  One codebase — Android, iOS, and Web",
    "✅  150+ tests — reliable, rule-accurate engine",
    "✅  Open-source stack — zero licensing cost",
]
for i, msg in enumerate(msgs):
    add_text(sl, msg, Inches(0.5), Inches(3.3 + i * 0.55),
             Inches(4.6), Inches(0.48), size=11, color=WHITE)

# Right — stats summary
add_text(sl, "By the Numbers",
         Inches(5.85), Inches(0.8), Inches(6.8), Inches(0.65),
         size=22, bold=True, color=GREEN)

summary_stats = [
    ("5",     "Bottom navigation tabs"),
    ("4",     "Match formats supported"),
    ("12",    "Dismissal types handled"),
    ("150+",  "Engine unit tests"),
    ("50 mi", "Live score broadcast radius"),
    ("5",     "Cloud Supabase tables"),
    ("3",     "Supported platforms"),
    ("0",     "Licensing cost (open-source stack)"),
]
for i, (val, lbl) in enumerate(summary_stats):
    col = i % 2; row = i // 2
    bx = Inches(5.85 + col * 3.5)
    cy = Inches(1.55 + row * 1.3)
    add_rect(sl, bx, cy, Inches(3.2), Inches(1.1),
             fill=GREEN_LIGHT, line=RGBColor(0xC0, 0xDC, 0xC2), line_width=Pt(0.5))
    add_text(sl, val, bx + Inches(0.15), cy + Inches(0.06),
             Inches(2.9), Inches(0.52), size=26, bold=True, color=GREEN,
             align=PP_ALIGN.CENTER)
    add_text(sl, lbl, bx + Inches(0.1), cy + Inches(0.6),
             Inches(3.0), Inches(0.38), size=9, color=GREY, align=PP_ALIGN.CENTER)

add_text(sl, "docs/ARCHITECTURE.md — full design diagrams",
         Inches(5.85), Inches(7.0), Inches(6.8), Inches(0.38),
         size=9, color=GREY, italic=True)


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), '..', 'docs', 'Inningsly_Executive.pptx')
prs.save(out_path)
print(f"✅  Saved → {os.path.abspath(out_path)}")
